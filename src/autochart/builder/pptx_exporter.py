"""Export AutoChart data to a branded BPHC PowerPoint presentation.

Creates one slide per chart with:
  - Title placeholder ("Insert title here")
  - Chart title (from data)
  - Native editable PowerPoint clustered bar chart
  - Data table below the chart
  - Footnote text box
  - Comment placeholder ("Insert comment here")

Uses the BPHC-branded .pptx template for slide master/theme.
"""

from __future__ import annotations

import io
import copy
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from lxml import etree

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "template_packages" / "bphc_template.pptx"

_TITLE_ONLY_LAYOUT = 5

# Colors — match the Excel chart scheme exactly
# The Excel bars use schemeClr tx2 + lumMod 25000 + lumOff 75000 = light blue
# In RGB that's approximately #B4C7E7
_BAR_COLOR = RGBColor(0xB4, 0xC7, 0xE7)
_BAR_COLOR_HEX = "B4C7E7"

# Set A uses distinct series colors
_GREEN = RGBColor(0x92, 0xD0, 0x50)
_BLUE = RGBColor(0x00, 0x70, 0xC0)
_NAVY = RGBColor(0x0E, 0x28, 0x41)

_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_BLACK = RGBColor(0x00, 0x00, 0x00)
_GRAY = RGBColor(0x59, 0x56, 0x59)

# Positions (EMU)
_CHART_TITLE_LEFT = Emu(1097280)
_CHART_TITLE_TOP = Emu(1740393)
_CHART_TITLE_H = Emu(369332)

_CHART_LEFT = Emu(410564)
_CHART_TOP = Emu(2150000)
_CHART_W = Emu(6400000)
_CHART_H = Emu(3900000)

_TABLE_LEFT = Emu(410564)
_TABLE_TOP = Emu(5950000)
_TABLE_H = Emu(600000)
_SLIDE_BOTTOM = Emu(6700000)  # leave margin above the gold bar at 6858000

_COMMENT_LEFT = Emu(7266547)
_COMMENT_TOP = Emu(2315431)
_COMMENT_W = Emu(4663063)
_COMMENT_H = Emu(2500000)

_FOOTNOTE_LEFT = Emu(7019999)
_FOOTNOTE_TOP = Emu(5200000)
_FOOTNOTE_W = Emu(5046688)
_FOOTNOTE_H = Emu(1200000)


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class ChartSeries:
    """A single data series in a chart."""
    name: str
    values: list[float]
    color: RGBColor | None = None
    point_colors: dict[int, RGBColor] = field(default_factory=dict)


@dataclass
class SlideData:
    """All data needed to create one slide."""
    chart_title: str
    categories: list[str]
    series: list[ChartSeries]
    footnote_lines: list[str]
    description: str = ""
    rate_unit: str = "per 100,000 residents"
    asterisk_points: list[tuple[int, int]] = field(default_factory=list)
    # Indices of bars that should get diagonal stripe pattern
    pattern_points: list[tuple[int, int]] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _add_text_box(
    slide, left, top, width, height, text,
    font_name="Montserrat", font_size=Pt(11),
    bold=False, color=_BLACK, alignment=PP_ALIGN.LEFT,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_multi_line_text_box(
    slide, left, top, width, height, lines,
    font_name="Montserrat", font_size=Pt(10), color=_BLACK,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = color
        p.space_after = Pt(2)
    return txBox


# ---------------------------------------------------------------------------
# Data table — clean style matching BPHC aesthetic
# ---------------------------------------------------------------------------

def _add_data_table(slide, left, top, width, categories, series_list, series_colors):
    """Add a clean data table below the chart.

    Layout matches the Excel output: colored legend square + series name in
    first column, then one column per category with centered values.
    No heavy header styling — just clean lines.
    """
    rows = len(series_list) + 1  # header + data
    cols = 1 + len(categories)

    # Compute row height so the table fits above the slide bottom
    max_table_h = _SLIDE_BOTTOM - top
    row_h = min(Emu(220000), max_table_h // rows)
    table_h = row_h * rows

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, table_h)
    table = table_shape.table

    # Disable built-in banding/styling (removes default blue/purple theme)
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is None:
        tblPr = etree.SubElement(tbl, qn('a:tblPr'))
    tblPr.set('bandRow', '0')
    tblPr.set('bandCol', '0')
    tblPr.set('firstRow', '0')
    tblPr.set('lastRow', '0')
    # Remove any tblStyle that forces theme colors
    for style_el in tblPr.findall(qn('a:tblStyle')):
        tblPr.remove(style_el)
    # Clear the tblStyle attribute if present
    if tblPr.get('tblStyle'):
        del tblPr.attrib['tblStyle']

    def _style_cell(cell, text, font_sz=Pt(8), bold=False, align=PP_ALIGN.CENTER,
                    font_color=_BLACK):
        cell.text = text
        for p in cell.text_frame.paragraphs:
            p.font.name = "Montserrat"
            p.font.size = font_sz
            p.font.bold = bold
            p.font.color.rgb = font_color
            p.alignment = align
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        # White background
        tcPr = cell._tc.get_or_add_tcPr()
        # Remove any existing fills
        for existing in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(existing)
        for existing in tcPr.findall(qn('a:noFill')):
            tcPr.remove(existing)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        etree.SubElement(solidFill, qn('a:srgbClr'), val='FFFFFF')
        # Thin gray borders
        for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            for existing in tcPr.findall(qn(border_name)):
                tcPr.remove(existing)
            ln = etree.SubElement(tcPr, qn(border_name), w='6350')
            lnFill = etree.SubElement(ln, qn('a:solidFill'))
            etree.SubElement(lnFill, qn('a:srgbClr'), val='D9D9D9')

    # Header row
    _style_cell(table.cell(0, 0), "", bold=True, align=PP_ALIGN.LEFT)
    for j, cat in enumerate(categories):
        _style_cell(table.cell(0, j + 1), cat, bold=True)

    # Data rows
    for i, s in enumerate(series_list):
        _style_cell(table.cell(i + 1, 0), s.name, bold=True, align=PP_ALIGN.LEFT)
        for j, val in enumerate(s.values):
            _style_cell(table.cell(i + 1, j + 1), f"{val:.1f}")

    return table_shape


# ---------------------------------------------------------------------------
# Pattern fill via OOXML post-processing
# ---------------------------------------------------------------------------

def _apply_pattern_fill_to_chart(chart, series_idx, point_idx):
    """Apply diagonal stripe pattern to a specific data point in a chart.

    python-pptx doesn't support pattFill, so we manipulate the XML directly.
    """
    plot_el = chart._chartSpace.findall('.//' + qn('c:barChart'))
    if not plot_el:
        return
    plot_el = plot_el[0]

    ser_els = plot_el.findall(qn('c:ser'))
    if series_idx >= len(ser_els):
        return
    ser = ser_els[series_idx]

    # Find or create dPt for this point
    dPt = None
    for existing in ser.findall(qn('c:dPt')):
        idx_el = existing.find(qn('c:idx'))
        if idx_el is not None and idx_el.get('val') == str(point_idx):
            dPt = existing
            break

    if dPt is None:
        dPt = etree.SubElement(ser, qn('c:dPt'))
        idx_el = etree.SubElement(dPt, qn('c:idx'))
        idx_el.set('val', str(point_idx))

    # Get or create spPr
    spPr = dPt.find(qn('c:spPr'))
    if spPr is None:
        spPr = etree.SubElement(dPt, qn('c:spPr'))

    # Remove existing fills
    for tag in [qn('a:solidFill'), qn('a:pattFill')]:
        for existing in spPr.findall(tag):
            spPr.remove(existing)

    # Add pattern fill matching the Excel format exactly
    pattFill = etree.SubElement(spPr, qn('a:pattFill'))
    pattFill.set('prst', 'wdDnDiag')

    fgClr = etree.SubElement(pattFill, qn('a:fgClr'))
    fgScheme = etree.SubElement(fgClr, qn('a:srgbClr'))
    fgScheme.set('val', _BAR_COLOR_HEX)

    bgClr = etree.SubElement(pattFill, qn('a:bgClr'))
    bgScheme = etree.SubElement(bgClr, qn('a:srgbClr'))
    bgScheme.set('val', 'FFFFFF')


# ---------------------------------------------------------------------------
# Chart builder
# ---------------------------------------------------------------------------

def _add_chart(slide, left, top, width, height, slide_data):
    """Add a native editable clustered bar chart matching Excel aesthetic."""
    chart_data = CategoryChartData()
    chart_data.categories = slide_data.categories
    for s in slide_data.series:
        chart_data.add_series(s.name, s.values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data,
    )
    chart = chart_frame.chart

    # Legend: only for multi-series (Set A)
    chart.has_legend = len(slide_data.series) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.name = "Montserrat"
        chart.legend.font.size = Pt(8)

    chart.has_title = False

    # Plot styling
    plot = chart.plots[0]
    plot.gap_width = 150
    plot.overlap = -27

    # Data labels
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.name = "Montserrat"
    dl.font.size = Pt(9)
    dl.font.color.rgb = _GRAY
    dl.number_format = '0.0'
    dl.show_value = True
    dl.show_category_name = False
    dl.label_position = XL_LABEL_POSITION.OUTSIDE_END

    # Color the series
    for i, s in enumerate(slide_data.series):
        series = chart.series[i]
        if s.color:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = s.color
        # Remove series border/line (clean look)
        series.format.line.fill.background()

        for pt_idx, pt_color in s.point_colors.items():
            if pt_idx < len(s.values):
                point = series.points[pt_idx]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = pt_color

    # Apply pattern fills (diagonal stripes on White/reference bars)
    for ser_idx, pt_idx in slide_data.pattern_points:
        _apply_pattern_fill_to_chart(chart, ser_idx, pt_idx)

    # Axes styling
    cat_axis = chart.category_axis
    cat_axis.has_title = False
    cat_axis.tick_labels.font.name = "Montserrat"
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.color.rgb = _GRAY
    cat_axis.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
    cat_axis.has_major_gridlines = False

    val_axis = chart.value_axis
    val_axis.tick_labels.font.name = "Montserrat"
    val_axis.tick_labels.font.size = Pt(9)
    val_axis.tick_labels.font.color.rgb = _GRAY
    val_axis.tick_labels.number_format = '0.0'
    val_axis.format.line.fill.background()  # hide axis line
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
    val_axis.major_gridlines.format.line.width = Pt(0.5)

    # Y-axis title
    val_axis.has_title = True
    val_axis.axis_title.text_frame.paragraphs[0].text = _rate_unit_to_axis_title(
        slide_data.rate_unit
    )
    val_axis.axis_title.text_frame.paragraphs[0].font.name = "Montserrat"
    val_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(8)
    val_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = _GRAY

    # Remove plot area border
    plot_area = chart._chartSpace.findall(qn('c:chart') + '/' + qn('c:plotArea'))
    if not plot_area:
        plot_area = [el for el in chart._chartSpace.iter(qn('c:plotArea'))]
    if plot_area:
        spPr = plot_area[0].find(qn('c:spPr'))
        if spPr is None:
            spPr = etree.SubElement(plot_area[0], qn('c:spPr'))
        for existing in spPr.findall(qn('a:ln')):
            spPr.remove(existing)
        ln = etree.SubElement(spPr, qn('a:ln'))
        etree.SubElement(ln, qn('a:noFill'))

    return chart_frame


def _rate_unit_to_axis_title(rate_unit: str) -> str:
    """Convert rate unit string to Y-axis title."""
    # "per 100,000 residents" -> "Rate per 100,000 Residents"
    if rate_unit:
        return f"Rate {rate_unit.strip()}"
    return "Rate"


# ---------------------------------------------------------------------------
# Slide builder
# ---------------------------------------------------------------------------

def _build_slide(prs, slide_data):
    layout = prs.slide_layouts[_TITLE_ONLY_LAYOUT]
    slide = prs.slides.add_slide(layout)

    # 1. Title placeholder
    if slide.placeholders:
        title_ph = slide.placeholders[0]
        title_ph.text = "Insert title here"
        for p in title_ph.text_frame.paragraphs:
            p.font.name = "Montserrat"
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = _NAVY

    # 2. Chart title
    _add_text_box(
        slide, _CHART_TITLE_LEFT, _CHART_TITLE_TOP,
        _CHART_W + Emu(2000000), _CHART_TITLE_H,
        slide_data.chart_title,
        font_name="Montserrat", font_size=Pt(16), color=_BLACK,
    )

    # 3. Chart
    _add_chart(slide, _CHART_LEFT, _CHART_TOP, _CHART_W, _CHART_H, slide_data)

    # 4. Data table
    series_colors = [s.color for s in slide_data.series]
    _add_data_table(
        slide, _TABLE_LEFT, _TABLE_TOP, _CHART_W,
        slide_data.categories, slide_data.series, series_colors,
    )

    # 5. Comment placeholder
    _add_text_box(
        slide, _COMMENT_LEFT, _COMMENT_TOP, _COMMENT_W, _COMMENT_H,
        "Insert comment here",
        font_name="Montserrat", font_size=Pt(14), color=_GRAY,
    )

    # 6. Footnote
    if slide_data.footnote_lines:
        _add_multi_line_text_box(
            slide, _FOOTNOTE_LEFT, _FOOTNOTE_TOP, _FOOTNOTE_W, _FOOTNOTE_H,
            slide_data.footnote_lines,
            font_name="Montserrat", font_size=Pt(10), color=_BLACK,
        )


# ---------------------------------------------------------------------------
# Data conversion: AutoChart data → SlideData
# ---------------------------------------------------------------------------

def _slides_from_set_a(data_list, config):
    from autochart.text.generator import TextGenerator
    from autochart.config import ChartSetType
    gen = TextGenerator(config)
    slides = []
    for d in data_list:
        title = gen.chart_title(ChartSetType.A, d.race_name)
        slides.append(SlideData(
            chart_title=title,
            categories=["Boston", "Female", "Male"],
            series=[
                ChartSeries(name=d.race_name, color=_GREEN,
                            values=[d.boston.group_rate, d.female.group_rate, d.male.group_rate]),
                ChartSeries(name="Rest of Boston", color=_BLUE,
                            values=[d.boston.reference_rate, d.female.reference_rate, d.male.reference_rate]),
                ChartSeries(name="Boston Overall", color=_NAVY,
                            values=[d.boston_overall_rate, d.female_overall_rate, d.male_overall_rate]),
            ],
            footnote_lines=gen.footnote().split("\n"),
            description=gen.descriptive_text_set_a(d),
            rate_unit=config.rate_unit,
        ))
    return slides


def _slides_from_set_b(data_list, config):
    from autochart.text.generator import TextGenerator
    from autochart.config import ChartSetType
    gen = TextGenerator(config)
    slides = []
    for d in data_list:
        title = gen.chart_title(ChartSetType.B, d.race_name)
        asterisks = [(0, 0)] if d.comparison.is_significant else []
        slides.append(SlideData(
            chart_title=title,
            categories=[d.race_name, config.reference_group, config.geography],
            series=[ChartSeries(
                name="Rate",
                values=[d.comparison.group_rate, d.comparison.reference_rate, d.boston_overall_rate],
                color=_BAR_COLOR,
            )],
            footnote_lines=gen.footnote().split("\n"),
            description=gen.descriptive_text_set_b(d),
            rate_unit=config.rate_unit,
            asterisk_points=asterisks,
            pattern_points=[(0, 1)],  # White bar = index 1
        ))
    return slides


def _slides_from_set_c(data_list, config):
    from autochart.text.generator import TextGenerator
    from autochart.config import ChartSetType
    gen = TextGenerator(config)
    slides = []
    for d in data_list:
        title = gen.chart_title(ChartSetType.C)
        categories = [c.group_name for c in d.comparisons] + [config.reference_group, config.geography]
        values = [c.group_rate for c in d.comparisons] + [d.comparisons[0].reference_rate, d.boston_overall_rate]
        white_idx = len(d.comparisons)
        asterisks = [(0, i) for i, c in enumerate(d.comparisons) if c.is_significant]

        slides.append(SlideData(
            chart_title=title,
            categories=categories,
            series=[ChartSeries(name="Rate", values=values, color=_BAR_COLOR)],
            footnote_lines=gen.footnote().split("\n"),
            description=gen.descriptive_text_set_c(d),
            rate_unit=config.rate_unit,
            asterisk_points=asterisks,
            pattern_points=[(0, white_idx)],
        ))
    return slides


def _slides_from_part3(data_list, config):
    from autochart.text.generator import TextGenerator
    from autochart.config import ChartSetType
    gen = TextGenerator(config)
    slides = []
    for d in data_list:
        title = gen.chart_title(ChartSetType.PART_3)
        f_names = [c.group_name for c in d.female_comparisons]
        categories = (
            [f"F: {n}" for n in f_names] + [f"F: {config.reference_group}", f"F: {config.geography}"]
            + [f"M: {n}" for n in f_names] + [f"M: {config.reference_group}", f"M: {config.geography}"]
        )
        f_vals = [c.group_rate for c in d.female_comparisons] + [d.female_comparisons[0].reference_rate, d.female_boston_rate]
        m_vals = [c.group_rate for c in d.male_comparisons] + [d.male_comparisons[0].reference_rate, d.male_boston_rate]
        white_f = len(d.female_comparisons)
        white_m = len(f_vals) + len(d.male_comparisons)

        slides.append(SlideData(
            chart_title=title,
            categories=categories,
            series=[ChartSeries(name="Rate", values=f_vals + m_vals, color=_BAR_COLOR)],
            footnote_lines=gen.footnote().split("\n"),
            description=gen.descriptive_text_part3(d),
            rate_unit=config.rate_unit,
            pattern_points=[(0, white_f), (0, white_m)],
        ))
    return slides


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def export_to_pptx(sheet_results, template_path=None):
    from autochart.config import ChartSetType

    if template_path is None:
        template_path = _TEMPLATE_PATH

    prs = Presentation(str(template_path))

    # Remove existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    all_slides = []
    seen = set()
    for sr in sheet_results:
        for ct, data_list in sr.by_type.items():
            key = (sr.config.disease_name, ct.value)
            if key in seen:
                continue
            seen.add(key)
            if ct == ChartSetType.A:
                all_slides.extend(_slides_from_set_a(data_list, sr.config))
            elif ct == ChartSetType.B:
                all_slides.extend(_slides_from_set_b(data_list, sr.config))
            elif ct == ChartSetType.C:
                all_slides.extend(_slides_from_set_c(data_list, sr.config))
            elif ct == ChartSetType.PART_3:
                all_slides.extend(_slides_from_part3(data_list, sr.config))

    for sd in all_slides:
        _build_slide(prs, sd)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
